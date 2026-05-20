# ================================================================
#  محلل حسابات X الاستخباراتي
#  الإصدار: v10.3.2
#  التحديثات:
#   - إضافة ميزة التحليل الاستخباراتي الشامل للمنشورات
#   - دعم جلب آخر 500 تغريدة عبر TwitterAPI.io
#   - تحليل الكلمات الدلالية المعادية للمملكة العربية السعودية
#   - تقرير مستوى الخطورة والتصنيف الاستخباراتي
#   - RTL كامل لجميع عناصر الواجهة
#   - إزالة Nitter - الاعتماد على FxTwitter API
# ================================================================

import streamlit as st
import requests
import base64
import io
import re
import random
import time
from datetime import datetime
from PIL import Image

try:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor, Emu
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PInches, Pt as PPt, Emu as PEmu
    from pptx.dml.color import RGBColor as PRGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import google.generativeai as genai
    GENAI_AVAILABLE = True
except ImportError:
    GENAI_AVAILABLE = False

# ════════════════════════════════════════════════════
#  الثوابت
# ════════════════════════════════════════════════════

PAGE_TITLE = "محلل حسابات X الاستخباراتي"
PAGE_ICON  = "🔍"
VERSION    = "v10.3.2"

USER_AGENTS = [
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
    "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36",
    "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
    "AppleWebKit/605.1.15 (KHTML, like Gecko) Version/17.4 Safari/605.1.15",
    "Mozilla/5.0 (X11; Linux x86_64; rv:125.0) Gecko/20100101 Firefox/125.0",
]

FXTWITTER_TWEET = "https://api.fxtwitter.com/status/{tweet_id}"
FXTWITTER_USER  = "https://api.fxtwitter.com/{username}"
TWITTERAPI_BASE = "https://api.twitterapi.io"

GEMINI_MODELS = {
    "Gemini 2.5 Flash (موصى به)":       "gemini-2.5-flash",
    "Gemini 2.0 Flash Lite (اقتصادي)":  "gemini-2.0-flash-lite",
    "Gemini 2.5 Pro (متقدم)":            "gemini-2.5-pro",
    "Gemini 1.5 Flash (احتياطي)":        "gemini-1.5-flash",
}

# ── الكلمات الدلالية الاستخباراتية ──
INTEL_KEYWORDS = {
    "🔴 عداء مباشر للمملكة": [
        "ارض الحرمين", "المهلكة", "ال سلول", "آل سلول",
        "شولوم", "البقرة", "الحلوب", "العلوج",
        "الغنيمة", "سلاليم", "النظام السعودي"
    ],
    "🟠 تحريض سياسي": [
        "قمع", "فساد", "إسقاط النظام", "تكميم الأفواه",
        "اعتقالات سياسية", "معتقل رأي", "سجن ظلم",
        "انتهاك حقوق", "سجناء الرأي", "بطش الحكومة"
    ],
    "🟡 حملات إعلامية عدائية": [
        "مقاطعة السعودية", "حملة ضد السعودية",
        "فضائح السعودية", "فشل سعودي", "الذباب الإلكتروني",
        "غسيل سمعة", "استهداف السعودية", "زعزعة الأمن",
        "إسقاط الحكم", "التحريض ضد الدولة", "الفوضى",
        "التجنيد ضد السعودية", "البطالة", "العاطلين"
    ],
    "🟣 إساءة للهوية الوطنية": [
        "كراهية السعوديين", "العنصرية ضد السعوديين",
        "إهانة رموز الدولة", "التشكيك بالوطنية",
        "الإساءة للهوية الوطنية", "ازدراء السعودية"
    ],
    "⚫ تحريض على العصيان": [
        "تمرد", "عصيان", "انتفاضة", "الشارع يغلي",
        "ثورة", "الدعوة للتظاهر", "مظاهرة",
        "قمع حريات", "انفجار شعبي", "حراك شعبي",
        "الغضب الشعبي", "العصيان المدني", "الربيع العربي"
    ],
}

IMAGE_ANALYSIS_POINTS = {
    "📍 تحديد الموقع الجغرافي": (
        "حلل هذه الصورة لتحديد الموقع الجغرافي. ابحث عن: "
        "المعالم البارزة، اللافتات، الطرازات المعمارية، الغطاء النباتي، "
        "التضاريس، المركبات، ملابس السكان. "
        "قدم احتمالات متعددة للمواقع مع مستويات ثقتك."
    ),
    "👤 تحليل الأشخاص والهويات": (
        "حلل الأشخاص الظاهرين. صف: السن، الجنس، الملابس، "
        "الشعارات، لغة الجسد، العلاقات، المعدات."
    ),
    "🚗 تحليل المركبات والمعدات": (
        "حدد وحلل المركبات والمعدات. اذكر: النوع، الماركة، "
        "الحالة، أرقام اللوحات، التحصينات، الاستخدام المحتمل."
    ),
    "📄 تحليل الوثائق والنصوص": (
        "اقرأ وحلل النصوص والوثائق. استخرج: النصوص، أرقام الوثائق، "
        "التواريخ، الأختام، الرموز الرسمية."
    ),
    "⚔️ تحليل الأسلحة والتجهيزات": (
        "حلل الأسلحة والتجهيزات. صف: النوع، الطراز، الحالة، "
        "طريقة الحمل، التكتيكات، مستوى التدريب."
    ),
    "🏗️ تحليل البنية التحتية": (
        "حلل البنية التحتية والمنشآت. صف: النوع، الحالة، "
        "الأهمية الاستراتيجية، علامات الإنشاء أو الهدم."
    ),
    "⏰ تحليل التوقيت والزمن": (
        "حلل مؤشرات التوقيت. ابحث عن: موضع الشمس، الظلال، "
        "الطقس، الساعات أو التقويمات، علامات القِدَم أو الحداثة."
    ),
    "🔍 تحليل الأحداث والتجمعات": (
        "حلل الحدث أو التجمع. صف: الطبيعة، حجم الحشد، "
        "التنظيم، المزاج، الشعارات، مستوى الأمن."
    ),
    "🔎 كشف التزوير والتلاعب": (
        "افحص الصورة لكشف التلاعب. ابحث عن: عدم التناسق في الإضاءة، "
        "حواف غير طبيعية، تكرار البكسل، علامات الذكاء الاصطناعي. "
        "قدم تقييماً لمصداقية الصورة."
    ),
    "🧠 تحليل شامل ومتكامل": (
        "قم بتحليل استخباراتي شامل. اجمع: الموقع، الأشخاص، "
        "المركبات، الأسلحة، البنية التحتية، التوقيت، السياق."
    ),
}


# ════════════════════════════════════════════════════
#  دوال مساعدة
# ════════════════════════════════════════════════════

def safe_text(val, default="غير متوفر"):
    if val is None:
        return default
    s = str(val).strip()
    return s if s else default


def escape_html(text: str) -> str:
    return (str(text)
            .replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;"))


def nl_to_br(text: str) -> str:
    return escape_html(text).replace("\n", "<br>")


def extract_tweet_id(raw: str):
    raw = raw.strip()
    m = re.search(r"/status(?:es)?/(\d+)", raw)
    if m:
        return m.group(1)
    if raw.isdigit() and len(raw) > 5:
        return raw
    return None


def extract_username(raw: str) -> str:
    raw = raw.strip().lstrip("@")
    m = re.search(r"(?:twitter\.com|x\.com)/([^/?#\s]+)", raw)
    if m:
        c = m.group(1)
        if c.lower() not in ("status", "i", "intent", "search", "hashtag"):
            return c
    m = re.match(r"^([A-Za-z0-9_]{1,50})$", raw)
    if m:
        return m.group(1)
    return raw


def format_number(n) -> str:
    try:
        n = int(n)
        if n >= 1_000_000:
            return f"{n/1_000_000:.1f}M"
        if n >= 1_000:
            return f"{n/1_000:.1f}K"
        return str(n)
    except Exception:
        return str(n)


def format_date(ts) -> str:
    try:
        return datetime.utcfromtimestamp(int(ts)).strftime("%Y-%m-%d %H:%M UTC")
    except Exception:
        return str(ts)


def image_to_base64(img_obj) -> str:
    buf = io.BytesIO()
    fmt = getattr(img_obj, "format", None) or "PNG"
    if fmt.upper() in ("JPEG", "JPG"):
        img_obj.save(buf, format="JPEG", quality=85)
    else:
        img_obj.save(buf, format="PNG")
    return base64.b64encode(buf.getvalue()).decode()


def base64_to_bytesio(b64_str: str) -> io.BytesIO:
    return io.BytesIO(base64.b64decode(b64_str))


def download_image_b64(url: str, timeout: int = 15):
    try:
        headers = {"User-Agent": random.choice(USER_AGENTS)}
        r = requests.get(url, headers=headers, timeout=timeout)
        r.raise_for_status()
        img = Image.open(io.BytesIO(r.content))
        return image_to_base64(img)
    except Exception as e:
        st.warning(f"تعذّر تحميل الصورة: {e}")
        return None


# ════════════════════════════════════════════════════
#  حساب أبعاد الصور
# ════════════════════════════════════════════════════

def calc_image_size_word(b64_str, max_w_inch=5.5, max_h_inch=7.0):
    try:
        bio = base64_to_bytesio(b64_str)
        img = Image.open(bio)
        w, h = img.size
        if w == 0 or h == 0:
            return {"width": Inches(4)}
        aspect = h / w
        if aspect * max_w_inch <= max_h_inch:
            return {"width": Inches(max_w_inch)}
        return {"height": Inches(max_h_inch)}
    except Exception:
        return {"width": Inches(4)}


def calc_image_size_pptx(b64_str, slide_w, slide_h,
                          max_w_ratio=0.78, max_h_ratio=0.62):
    try:
        bio = base64_to_bytesio(b64_str)
        img = Image.open(bio)
        w, h = img.size
        if w == 0 or h == 0:
            return (int(slide_w*0.1), int(slide_h*0.25),
                    int(slide_w*0.8), int(slide_h*0.5))
        aspect = h / w
        max_w  = int(slide_w * max_w_ratio)
        max_h  = int(slide_h * max_h_ratio)
        if aspect * max_w <= max_h:
            pic_w = max_w; pic_h = int(aspect * max_w)
        else:
            pic_h = max_h; pic_w = int(max_h / aspect)
        left = int((slide_w - pic_w) / 2)
        top  = PEmu(1_050_000) + int((max_h - pic_h) / 2)
        return left, top, pic_w, pic_h
    except Exception:
        return (int(slide_w*0.1), int(slide_h*0.25),
                int(slide_w*0.8), int(slide_h*0.5))


# ════════════════════════════════════════════════════
#  جلب البيانات - FxTwitter API
# ════════════════════════════════════════════════════

def fetch_fxtwitter_tweet(tweet_id: str):
    url = FXTWITTER_TWEET.format(tweet_id=tweet_id)
    headers = {"User-Agent": random.choice(USER_AGENTS), "Accept": "application/json"}
    try:
        r = requests.get(url, headers=headers, timeout=20)
        r.raise_for_status()
        data = r.json()
        if data.get("code") == 200 and "tweet" in data:
            return data["tweet"]
        return None
    except Exception as e:
        st.error(f"خطأ في جلب التغريدة: {e}")
        return None


def fetch_fxtwitter_user(username: str):
    username = username.strip().lstrip("@")
    url = FXTWITTER_USER.format(username=username)
    headers = {"User-Agent": random.choice(USER_AGENTS), "Accept": "application/json"}
    try:
        r = requests.get(url, headers=headers, timeout=20)
        r.raise_for_status()
        data = r.json()
        if data.get("code") == 200 and "user" in data:
            return data["user"]
        return None
    except Exception as e:
        st.error(f"خطأ في جلب بيانات الحساب: {e}")
        return None


# ════════════════════════════════════════════════════
#  جلب التغريدات - TwitterAPI.io
# ════════════════════════════════════════════════════

def fetch_user_tweets_twitterapi(username: str, tapi_key: str,
                                  limit: int = 500) -> list:
    """
    جلب آخر N تغريدة عبر TwitterAPI.io (مدفوع - رخيص جداً)
    يحتاج مفتاح من: https://twitterapi.io
    """
    username  = username.strip().lstrip("@")
    url       = f"{TWITTERAPI_BASE}/twitter/user/last_tweets"
    headers   = {
        "X-API-Key": tapi_key,
        "Content-Type": "application/json",
    }
    all_tweets = []
    cursor     = None
    max_pages  = 30  # حد أقصى للصفحات (20 تغريدة/صفحة = 600 تغريدة)

    progress = st.progress(0, text="جاري جلب التغريدات...")
    page = 0

    while len(all_tweets) < limit and page < max_pages:
        params = {"userName": username}
        if cursor:
            params["cursor"] = cursor
        try:
            r = requests.get(url, headers=headers, params=params, timeout=25)
            r.raise_for_status()
            data = r.json()

            tweets = data.get("tweets", [])
            if not tweets:
                break

            for t in tweets:
                text = (t.get("text") or
                        t.get("full_text") or
                        t.get("rawContent") or "")
                if text:
                    all_tweets.append(text.strip())

            cursor = (data.get("next_cursor") or
                      data.get("cursor")       or
                      data.get("nextCursor"))
            if not cursor:
                break

            page += 1
            pct = min(int((len(all_tweets) / limit) * 100), 99)
            progress.progress(pct,
                text=f"تم جلب {len(all_tweets)} تغريدة...")
            time.sleep(0.3)  # تجنب Rate Limit

        except requests.exceptions.HTTPError as e:
            if r.status_code == 401:
                st.error("❌ مفتاح TwitterAPI.io غير صحيح. تحقق منه.")
            elif r.status_code == 429:
                st.warning("⚠️ تجاوزت حد الطلبات. انتظر قليلاً وأعد المحاولة.")
            else:
                st.error(f"خطأ HTTP {r.status_code}: {e}")
            break
        except Exception as e:
            st.error(f"خطأ في جلب التغريدات: {e}")
            break

    progress.progress(100, text=f"✅ تم جلب {len(all_tweets)} تغريدة")
    time.sleep(0.5)
    progress.empty()

    return all_tweets[:limit]


# ════════════════════════════════════════════════════
#  الكشف اليدوي عن الكلمات الدلالية
# ════════════════════════════════════════════════════

def scan_keywords(tweets_text: str) -> dict:
    """فحص النص عن الكلمات الدلالية وإرجاع النتائج مصنّفة."""
    found = {}
    text_lower = tweets_text.lower()
    for category, keywords in INTEL_KEYWORDS.items():
        hits = []
        for kw in keywords:
            count = text_lower.count(kw.lower())
            if count > 0:
                hits.append((kw, count))
        if hits:
            found[category] = sorted(hits, key=lambda x: -x[1])
    return found


def threat_level_color(level: int) -> str:
    if level <= 2:
        return "#4caf50"   # أخضر
    if level <= 4:
        return "#8bc34a"   # أخضر فاتح
    if level <= 6:
        return "#ff9800"   # برتقالي
    if level <= 8:
        return "#f44336"   # أحمر
    return "#b71c1c"       # أحمر داكن


def threat_label(level: int) -> str:
    if level <= 2:
        return "آمن - لا توجد مؤشرات خطرة"
    if level <= 4:
        return "منخفض - مؤشرات بسيطة"
    if level <= 6:
        return "متوسط - محتوى مثير للقلق"
    if level <= 8:
        return "مرتفع - محتوى عدائي"
    return "خطر جداً - محتوى تحريضي"


# ════════════════════════════════════════════════════
#  Gemini
# ════════════════════════════════════════════════════

def _get_model(api_key: str, model_name: str):
    if not GENAI_AVAILABLE:
        st.error("مكتبة google-generativeai غير مثبّتة.")
        return None
    try:
        genai.configure(api_key=api_key)
        return genai.GenerativeModel(model_name)
    except Exception as e:
        st.error(f"خطأ في تهيئة Gemini: {e}")
        return None


def gemini_text(prompt: str, api_key: str, model_name: str):
    model = _get_model(api_key, model_name)
    if not model:
        return None
    try:
        resp = model.generate_content(prompt)
        return resp.text
    except Exception as e:
        st.error(f"خطأ في توليد النص: {e}")
        return None


def gemini_with_images(prompt: str, images_b64: list,
                        api_key: str, model_name: str):
    model = _get_model(api_key, model_name)
    if not model:
        return None
    try:
        parts = [prompt]
        for b64 in images_b64:
            raw  = base64.b64decode(b64)
            img  = Image.open(io.BytesIO(raw))
            fmt  = (img.format or "PNG").lower()
            mime = "image/jpeg" if fmt in ("jpeg","jpg") else "image/png"
            parts.append({"mime_type": mime, "data": b64})
        resp = model.generate_content(parts)
        return resp.text
    except Exception as e:
        st.error(f"خطأ في تحليل الصور: {e}")
        return None


def generate_intel_summary(tweets_list: list, profile_data: dict,
                            found_keywords: dict,
                            api_key: str, model_name: str) -> str:
    """توليد الملخص الاستخباراتي الشامل عبر Gemini."""

    # بناء قسم الكلمات المكتشفة
    kw_section = ""
    for cat, hits in found_keywords.items():
        kw_section += f"\n{cat}:\n"
        for kw, cnt in hits:
            kw_section += f"  - '{kw}' تكرر {cnt} مرة\n"

    # اختيار عيّنة من التغريدات للتحليل
    sample_size = min(len(tweets_list), 80)
    sample = tweets_list[:sample_size]
    tweets_sample = "\n---\n".join(sample)

    # بيانات الحساب
    bio = profile_data.get("description", "")
    name = profile_data.get("name", "")
    username = profile_data.get("screen_name", "")
    followers = format_number(profile_data.get("followers", 0))
    location = profile_data.get("location", "")

    prompt = f"""أنت محلل استخباراتي متخصص في رصد المحتوى الرقمي المعادي للمملكة العربية السعودية.

## بيانات الحساب:
- الاسم: {name} (@{username})
- الوصف: {bio}
- الموقع: {location}
- عدد المتابعين: {followers}
- عدد التغريدات المُحللة: {len(tweets_list)}

## الكلمات الدلالية المكتشفة تلقائياً:
{kw_section if kw_section else "لم تُكتشف كلمات دلالية صريحة"}

## عيّنة من المنشورات ({sample_size} منشور):
{tweets_sample}

---

## المطلوب: تقرير استخباراتي شامل يتضمن:

### 1. التصنيف الاستخباراتي النهائي
اختر من: [معادٍ للمملكة | متطرف | محرِّض | مثير للرأي العام | ناشط سياسي | عادي]

### 2. مستوى الخطورة
رقم من 1 إلى 10 مع تبرير واضح

### 3. المحاور العدائية المكتشفة
حدد من القائمة التالية ما ينطبق:
- عداء مباشر للمملكة (الألقاب المسيئة، الإساءة للقيادة)
- تحريض سياسي (الدعوة لإسقاط النظام، نشر الفوضى)
- حملات إعلامية عدائية (المشاركة في حملات ضد السعودية)
- إساءة للهوية الوطنية (كراهية السعوديين، إهانة الرموز)
- تحريض على العصيان (الدعوة للتظاهر، زعزعة الاستقرار)

### 4. تحليل نمط النشر
- التكرار والانتظام
- مواضيع التركيز الرئيسية
- الأسلوب (مباشر/غير مباشر/ساخر/توثيقي)
- مؤشرات الارتباط بمنظمات أو حملات منسقة

### 5. مؤشرات التطرف
تحديد أي أنماط تشير للتطرف الفكري أو الديني أو السياسي

### 6. ملخص الخطر الأمني
فقرة واحدة تلخص مستوى الخطورة على الأمن الوطني السعودي

### 7. التوصيات
- إجراءات فورية موصى بها
- مستوى الأولوية (عاجل/متوسط/منخفض)

قدم التقرير باللغة العربية بأسلوب رسمي واحترافي مناسب للجهات الأمنية."""

    return gemini_text(prompt, api_key, model_name)


# ════════════════════════════════════════════════════
#  مساعدات Word RTL
# ════════════════════════════════════════════════════

def _set_rtl_para(para):
    try:
        pPr  = para._p.get_or_add_pPr()
        bidi = OxmlElement("w:bidi")
        bidi.set(qn("w:val"), "1")
        pPr.append(bidi)
        jc = OxmlElement("w:jc")
        jc.set(qn("w:val"), "right")
        pPr.append(jc)
    except Exception:
        pass


def _add_word_section(doc, title, content, title_color=None, level=2):
    h = doc.add_heading(title, level=level)
    _set_rtl_para(h)
    if title_color:
        for run in h.runs:
            run.font.color.rgb = RGBColor(*title_color)
    p = doc.add_paragraph(content)
    _set_rtl_para(p)
    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    return p


# ════════════════════════════════════════════════════
#  مساعدات PPTX RTL
# ════════════════════════════════════════════════════

def _set_rtl_pptx(para):
    try:
        pPr  = para._p.get_or_add_pPr()
        bidi = OxmlElement("a:rtl")
        bidi.text = "1"
        pPr.append(bidi)
    except Exception:
        pass


# ════════════════════════════════════════════════════
#  تصدير Word
# ════════════════════════════════════════════════════

def export_to_word(title, tweet_data, img_analysis,
                   text_analysis, exec_summary, images_b64,
                   intel_report=""):
    if not DOCX_AVAILABLE:
        st.error("مكتبة python-docx غير مثبّتة.")
        return None
    doc   = Document()
    style = doc.styles["Normal"]
    style.font.name = "Arial"
    style.font.size = Pt(11)

    h0 = doc.add_heading(f"تقرير استخباراتي – {title}", 0)
    _set_rtl_para(h0)
    for run in h0.runs:
        run.font.color.rgb = RGBColor(0xE9, 0x45, 0x60)

    meta = doc.add_paragraph(
        f"التاريخ: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}"
        f"  |  الإصدار: {VERSION}"
    )
    _set_rtl_para(meta)
    doc.add_paragraph()

    # التقرير الاستخباراتي
    if intel_report:
        h_ir = doc.add_heading("🔴 التقرير الاستخباراتي الشامل للمنشورات", 1)
        _set_rtl_para(h_ir)
        for run in h_ir.runs:
            run.font.color.rgb = RGBColor(0xB7, 0x1C, 0x1C)
            run.font.size = Pt(16)
        p_ir = doc.add_paragraph(intel_report)
        _set_rtl_para(p_ir)
        p_ir.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        doc.add_paragraph()

    if exec_summary:
        h_es = doc.add_heading("📋 الملخص التنفيذي", 1)
        _set_rtl_para(h_es)
        for run in h_es.runs:
            run.font.color.rgb = RGBColor(0xE9, 0x45, 0x60)
            run.font.size = Pt(16)
        p_es = doc.add_paragraph(exec_summary)
        _set_rtl_para(p_es)
        p_es.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        doc.add_paragraph()

    if tweet_data:
        _add_word_section(
            doc, "📌 بيانات التغريدة",
            "\n".join(f"{k}: {v}" for k, v in tweet_data.items()),
            title_color=(0x1D, 0xA1, 0xF2),
        )
        doc.add_paragraph()

    if images_b64:
        h_img = doc.add_heading("🖼️ الصور المرفقة وتحليلها", 1)
        _set_rtl_para(h_img)
        for b64 in images_b64:
            try:
                size_kw = calc_image_size_word(b64)
                doc.add_picture(base64_to_bytesio(b64), **size_kw)
                doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
            except Exception as e:
                doc.add_paragraph(f"[تعذّر إدراج الصورة: {e}]")
        if img_analysis:
            _add_word_section(doc, "نتائج تحليل الصور", img_analysis)
        doc.add_paragraph()

    if text_analysis:
        _add_word_section(
            doc, "📝 تحليل نص التغريدة", text_analysis,
            title_color=(0x1D, 0xA1, 0xF2),
        )

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ════════════════════════════════════════════════════
#  تصدير PowerPoint
# ════════════════════════════════════════════════════

def export_to_pptx(title, tweet_data, img_analysis,
                    text_analysis, exec_summary, images_b64,
                    intel_report=""):
    if not PPTX_AVAILABLE:
        st.error("مكتبة python-pptx غير مثبّتة.")
        return None

    prs  = Presentation()
    sw   = prs.slide_width
    sh   = prs.slide_height
    DARK = PRGBColor(0x0D, 0x1B, 0x2A)
    RED  = PRGBColor(0xE9, 0x45, 0x60)
    DRED = PRGBColor(0xB7, 0x1C, 0x1C)
    BLUE = PRGBColor(0x1D, 0xA1, 0xF2)
    WHT  = PRGBColor(0xFF, 0xFF, 0xFF)
    GRAY = PRGBColor(0x88, 0x88, 0x88)

    def _new_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg    = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = DARK
        return slide

    def _add_tb(slide, text, x, y, w, h,
                color=None, size=16, bold=False, align=PP_ALIGN.RIGHT):
        color = color or WHT
        txb   = slide.shapes.add_textbox(x, y, w, h)
        tf    = txb.text_frame
        tf.word_wrap = True
        p     = tf.paragraphs[0]
        p.alignment = align
        _set_rtl_pptx(p)
        run   = p.add_run()
        run.text           = text
        run.font.color.rgb = color
        run.font.size      = PPt(size)
        run.font.bold      = bold

    def _red_line(slide, y_emu, color=RED):
        from pptx.util import Emu as _E
        shape = slide.shapes.add_shape(
            1, PInches(0.3), y_emu, sw - PInches(0.6), _E(45000)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    # ── شريحة العنوان ──
    s0 = _new_slide()
    _add_tb(s0, PAGE_TITLE, PInches(0.5), PInches(1.0),
            sw - PInches(1), PInches(1),
            color=RED, size=30, bold=True, align=PP_ALIGN.CENTER)
    _add_tb(s0, title, PInches(0.5), PInches(2.2),
            sw - PInches(1), PInches(0.8),
            color=WHT, size=22, align=PP_ALIGN.CENTER)
    _add_tb(s0, f"{datetime.utcnow().strftime('%Y-%m-%d')}  |  {VERSION}",
            PInches(0.5), PInches(3.2),
            sw - PInches(1), PInches(0.5),
            color=GRAY, size=14, align=PP_ALIGN.CENTER)

    # ── شريحة التقرير الاستخباراتي ──
    if intel_report:
        si = _new_slide()
        _add_tb(si, "🔴 التقرير الاستخباراتي الشامل",
                PInches(0.3), PInches(0.1),
                sw - PInches(0.6), PInches(0.6),
                color=DRED, size=22, bold=True)
        _red_line(si, PEmu(760_000), color=DRED)
        _add_tb(si, intel_report[:800],
                PInches(0.3), PInches(0.9),
                sw - PInches(0.6), sh - PInches(1.2),
                color=WHT, size=12)

    # ── شريحة الملخص التنفيذي ──
    if exec_summary:
        s_es = _new_slide()
        _add_tb(s_es, "📋 الملخص التنفيذي",
                PInches(0.3), PInches(0.1),
                sw - PInches(0.6), PInches(0.6),
                color=RED, size=22, bold=True)
        _red_line(s_es, PEmu(760_000))
        _add_tb(s_es, exec_summary[:700],
                PInches(0.3), PInches(0.9),
                sw - PInches(0.6), sh - PInches(1.2),
                color=WHT, size=13)

    # ── شرائح الصور ──
    for b64 in images_b64:
        si = _new_slide()
        _add_tb(si, "🖼️ تحليل الصورة",
                PInches(0.3), PInches(0.05),
                sw - PInches(0.6), PInches(0.5),
                color=BLUE, size=20, bold=True)
        left, top, pw, ph = calc_image_size_pptx(b64, sw, sh)
        try:
            si.shapes.add_picture(base64_to_bytesio(b64), left, top, pw, ph)
        except Exception as e:
            _add_tb(si, f"[تعذّر الصورة: {e}]",
                    PInches(0.3), PInches(1),
                    sw - PInches(0.6), PInches(1))

    # ── شريحة تحليل النص ──
    if text_analysis:
        st_txt = _new_slide()
        _add_tb(st_txt, "📝 تحليل النص",
                PInches(0.3), PInches(0.05),
                sw - PInches(0.6), PInches(0.5),
                color=BLUE, size=22, bold=True)
        _red_line(st_txt, PEmu(680_000))
        _add_tb(st_txt, text_analysis[:700],
                PInches(0.3), PInches(0.6),
                sw - PInches(0.6), sh - PInches(0.9),
                color=WHT, size=13)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ════════════════════════════════════════════════════
#  أزرار التصدير
# ════════════════════════════════════════════════════

def render_export_buttons(title: str = "report"):
    ss            = st.session_state
    tweet_data    = ss.get("tweet_data",    {})
    img_analysis  = ss.get("img_analysis",  "")
    text_analysis = ss.get("text_analysis", "")
    exec_summary  = ss.get("exec_summary",  "")
    images_b64    = ss.get("images_b64",    [])
    intel_report  = ss.get("intel_report",  "")

    if not any([tweet_data, img_analysis, text_analysis, intel_report]):
        return

    st.divider()
    st.subheader("📥 تصدير التقرير")
    c1, c2, c3 = st.columns(3)
    stamp = datetime.utcnow().strftime("%Y%m%d_%H%M")
    safe  = re.sub(r"[^\w\u0600-\u06FF]", "_", str(title))[:30]

    with c1:
        docx_bytes = export_to_word(
            title, tweet_data, img_analysis,
            text_analysis, exec_summary, images_b64, intel_report
        )
        if docx_bytes:
            st.download_button(
                "📄 Word", docx_bytes,
                file_name=f"X_Intel_{safe}_{stamp}.docx",
                mime=("application/vnd.openxmlformats-officedocument"
                      ".wordprocessingml.document"),
            )
    with c2:
        pptx_bytes = export_to_pptx(
            title, tweet_data, img_analysis,
            text_analysis, exec_summary, images_b64, intel_report
        )
        if pptx_bytes:
            st.download_button(
                "📊 PowerPoint", pptx_bytes,
                file_name=f"X_Intel_{safe}_{stamp}.pptx",
                mime=("application/vnd.openxmlformats-officedocument"
                      ".presentationml.presentation"),
            )
    with c3:
        parts = []
        if intel_report:
            parts.append(f"=== التقرير الاستخباراتي ===\n{intel_report}")
        if tweet_data:
            parts.append(
                "=== بيانات التغريدة ===\n"
                + "\n".join(f"{k}: {v}" for k, v in tweet_data.items())
            )
        if exec_summary:
            parts.append(f"=== الملخص التنفيذي ===\n{exec_summary}")
        if img_analysis:
            parts.append(f"=== تحليل الصور ===\n{img_analysis}")
        if text_analysis:
            parts.append(f"=== تحليل النص ===\n{text_analysis}")
        st.download_button(
            "📝 TXT", "\n\n".join(parts).encode("utf-8"),
            file_name=f"X_Intel_{safe}_{stamp}.txt",
            mime="text/plain",
        )


# ════════════════════════════════════════════════════
#  بطاقة عرض الحساب
# ════════════════════════════════════════════════════

def render_user_card(u: dict):
    name       = escape_html(safe_text(u.get("name")))
    screen     = escape_html(safe_text(u.get("screen_name")))
    desc       = nl_to_br(safe_text(u.get("description", ""), ""))
    location   = escape_html(safe_text(u.get("location", ""), ""))
    joined_raw = safe_text(u.get("joined", ""), "")
    joined     = joined_raw[:10] if joined_raw else ""
    followers  = format_number(u.get("followers",   0))
    following  = format_number(u.get("following",   0))
    tweets_cnt = format_number(u.get("tweets",      0))
    likes_cnt  = format_number(u.get("likes",       0))
    media_cnt  = format_number(u.get("media_count", 0))
    avatar     = u.get("avatar_url", "")
    banner     = u.get("banner_url", "")

    verif = u.get("verification", {})
    is_verified = (verif.get("verified", False)
                   if isinstance(verif, dict)
                   else bool(u.get("verified", False)))
    badge = "✅" if is_verified else ""

    avatar_hq   = avatar.replace("_normal", "_400x400") if avatar else ""
    avatar_html = (
        f'<img src="{escape_html(avatar_hq or avatar)}" '
        'style="width:90px;height:90px;border-radius:50%;'
        'border:3px solid #e94560;object-fit:cover;'
        'box-shadow:0 4px 12px rgba(233,69,96,0.4);">'
        if avatar else
        '<div style="width:90px;height:90px;border-radius:50%;'
        'background:#333;border:3px solid #e94560;"></div>'
    )
    banner_html = (
        f'<div style="background:url({escape_html(banner)}) center/cover;'
        'height:130px;border-radius:12px 12px 0 0;"></div>'
        if banner else
        '<div style="background:linear-gradient(135deg,#1a1a2e,#0d1b2a);'
        'height:70px;border-radius:12px 12px 0 0;"></div>'
    )

    loc_badge  = (
        f'<span style="background:#1a2a3a;padding:4px 12px;border-radius:20px;'
        f'color:#aaa;font-size:0.8rem;">📍 {location}</span>'
        if location else ""
    )
    join_badge = (
        f'<span style="background:#1a2a3a;padding:4px 12px;border-radius:20px;'
        f'color:#aaa;font-size:0.8rem;">📅 {joined}</span>'
        if joined else ""
    )

    st.markdown(f"""
    <div style="background:linear-gradient(135deg,#1a1a2e,#16213e);
                border-radius:14px;border:1px solid #e94560;
                margin-bottom:20px;overflow:hidden;
                box-shadow:0 4px 20px rgba(0,0,0,0.4);">
        {banner_html}
        <div style="padding:16px 20px 20px;">
            <div style="display:flex;align-items:flex-end;
                        gap:14px;margin-top:-45px;margin-bottom:14px;">
                {avatar_html}
                <div>
                    <div style="font-size:1.3rem;font-weight:700;
                                color:#fff;line-height:1.2;">
                        {name} {badge}
                    </div>
                    <div style="color:#1da1f2;font-size:1rem;">@{screen}</div>
                </div>
            </div>
            <div style="color:#ccc;font-size:0.92rem;direction:rtl;
                        line-height:1.7;margin-bottom:12px;">{desc}</div>
            <div style="display:flex;gap:8px;flex-wrap:wrap;
                        margin-bottom:14px;direction:rtl;">
                {loc_badge}{join_badge}
            </div>
            <div style="display:flex;gap:20px;flex-wrap:wrap;color:#ccc;
                        font-size:0.88rem;direction:rtl;
                        border-top:1px solid #2a3a4a;padding-top:12px;">
                <div><span style="color:#fff;font-weight:700;">{followers}</span> متابِع</div>
                <div><span style="color:#fff;font-weight:700;">{following}</span> يتابع</div>
                <div><span style="color:#fff;font-weight:700;">{tweets_cnt}</span> تغريدة</div>
                <div><span style="color:#fff;font-weight:700;">{likes_cnt}</span> إعجاب</div>
                <div><span style="color:#fff;font-weight:700;">{media_cnt}</span> وسائط</div>
            </div>
        </div>
    </div>
    """, unsafe_allow_html=True)


# ════════════════════════════════════════════════════
#  بطاقة عرض التغريدة
# ════════════════════════════════════════════════════

def render_tweet_card(tweet: dict):
    author  = tweet.get("author", {})
    name    = escape_html(safe_text(author.get("name")))
    screen  = escape_html(safe_text(author.get("screen_name")))
    text    = nl_to_br(safe_text(tweet.get("text", "")))
    likes   = format_number(tweet.get("likes",    0))
    rts     = format_number(tweet.get("retweets", 0))
    replies = format_number(tweet.get("replies",  0))
    views   = format_number(tweet.get("views",    0))
    ts      = format_date(tweet.get("created_timestamp", 0))
    url     = tweet.get("url", "#")

    st.markdown(f"""
    <div style="background:#16213e;border-radius:12px;padding:20px;
                border:1px solid #1da1f2;margin-bottom:16px;">
        <div style="display:flex;justify-content:space-between;
                    align-items:center;margin-bottom:12px;direction:rtl;">
            <div>
                <span style="color:#fff;font-weight:700;">{name}</span>
                <span style="color:#1da1f2;margin-right:8px;">@{screen}</span>
            </div>
            <span style="color:#888;font-size:0.82rem;">{ts}</span>
        </div>
        <div style="color:#e0e0e0;font-size:1rem;direction:rtl;
                    line-height:1.8;margin-bottom:16px;">{text}</div>
        <div style="display:flex;gap:24px;color:#888;font-size:0.9rem;direction:rtl;">
            <span>❤️ {likes}</span>
            <span>🔁 {rts}</span>
            <span>💬 {replies}</span>
            <span>👁️ {views}</span>
            <a href="{escape_html(url)}" target="_blank"
               style="color:#1da1f2;margin-right:auto;">🔗 فتح</a>
        </div>
    </div>
    """, unsafe_allow_html=True)


# ════════════════════════════════════════════════════
#  عرض نتائج الفحص الكلمي
# ════════════════════════════════════════════════════

def render_keyword_scan(found: dict, total_tweets: int):
    """عرض نتائج الفحص الكلمي بشكل بصري."""
    if not found:
        st.markdown("""
        <div style="background:#0a2a0a;border-radius:12px;padding:16px;
                    border:1px solid #4caf50;direction:rtl;text-align:right;">
            ✅ <strong style="color:#4caf50;">لم تُكتشف كلمات دلالية عدائية</strong>
            في عيّنة المنشورات المُحللة.
        </div>
        """, unsafe_allow_html=True)
        return

    total_hits = sum(sum(c for _, c in hits) for hits in found.values())

    # بطاقة ملخص
    st.markdown(f"""
    <div style="background:linear-gradient(135deg,#2a0a0a,#1a0510);
                border-radius:12px;padding:16px 20px;
                border:2px solid #e94560;direction:rtl;
                margin-bottom:16px;">
        <div style="display:flex;justify-content:space-between;align-items:center;">
            <div>
                <div style="color:#e94560;font-size:1.1rem;font-weight:700;">
                    ⚠️ تم اكتشاف كلمات دلالية عدائية
                </div>
                <div style="color:#ccc;font-size:0.9rem;margin-top:4px;">
                    إجمالي التكرارات: <strong style="color:#ff9800;">{total_hits}</strong>
                    في <strong style="color:#fff;">{total_tweets}</strong> منشور
                </div>
            </div>
            <div style="font-size:2rem;">{len(found)}</div>
        </div>
    </div>
    """, unsafe_allow_html=True)

    # تفاصيل حسب الفئة
    for cat, hits in found.items():
        cat_total = sum(c for _, c in hits)
        kw_html   = " ".join(
            f'<span style="background:#2a1a1a;border:1px solid #e94560;'
            f'border-radius:20px;padding:3px 10px;font-size:0.82rem;'
            f'color:#ff9800;">"{kw}" ×{cnt}</span>'
            for kw, cnt in hits
        )
        st.markdown(f"""
        <div style="background:#1a1a2e;border-radius:10px;padding:14px;
                    border-right:4px solid #e94560;margin-bottom:10px;direction:rtl;">
            <div style="display:flex;justify-content:space-between;
                        align-items:center;margin-bottom:8px;">
                <strong style="color:#fff;">{cat}</strong>
                <span style="background:#e94560;color:#fff;padding:2px 10px;
                             border-radius:12px;font-size:0.8rem;">
                    {cat_total} تكرار
                </span>
            </div>
            <div style="display:flex;flex-wrap:wrap;gap:6px;">
                {kw_html}
            </div>
        </div>
        """, unsafe_allow_html=True)


# ════════════════════════════════════════════════════
#  تبويب تحليل الحساب
# ════════════════════════════════════════════════════

def account_tab(api_key: str, model_name: str):
    st.header("🔍 تحليل حساب X")

    # ── جلب بيانات الحساب ──
    col1, col2 = st.columns([3, 1])
    with col1:
        raw_input = st.text_input(
            "رابط الحساب أو اسم المستخدم",
            placeholder="مثال:  https://x.com/username  أو  username",
            key="account_input",
        )
    with col2:
        uploaded_img = st.file_uploader(
            "صورة (اختياري)",
            type=["jpg", "png", "webp"],
            key="account_img_upload",
        )

    if st.button("🔎 جلب بيانات الحساب", key="btn_fetch_account"):
        if not raw_input.strip():
            st.error("يرجى إدخال رابط الحساب أو اسم المستخدم.")
            return

        tweet_id = extract_tweet_id(raw_input)

        if tweet_id:
            with st.spinner("⏳ جاري جلب بيانات التغريدة..."):
                tweet = fetch_fxtwitter_tweet(tweet_id)
            if not tweet:
                st.error("❌ تعذّر جلب بيانات التغريدة.")
                return
            uname = tweet.get("author", {}).get("screen_name", "")
            with st.spinner(f"⏳ جاري جلب ملف الحساب @{uname}..."):
                user_data = fetch_fxtwitter_user(uname) if uname else None
            if not user_data:
                author = tweet.get("author", {})
                user_data = {
                    "screen_name": author.get("screen_name", ""),
                    "name":        author.get("name", ""),
                    "avatar_url":  author.get("avatar_url", ""),
                    "banner_url":  author.get("banner_url", ""),
                }
            st.session_state["account_user"]  = user_data
            st.session_state["account_tweet"] = tweet
        else:
            username = extract_username(raw_input)
            with st.spinner(f"⏳ جاري جلب بيانات @{username}..."):
                user_data = fetch_fxtwitter_user(username)
            if not user_data:
                st.error(
                    "❌ تعذّر جلب بيانات الحساب.\n\n"
                    "تأكد من صحة الاسم، أو جرّب رابط تغريدة مباشرة."
                )
                return
            st.session_state["account_user"] = user_data
            st.session_state.pop("account_tweet", None)

        # مسح التقارير القديمة
        for key in ["intel_tweets", "intel_report", "img_analysis"]:
            st.session_state.pop(key, None)
        st.success("✅ تم جلب بيانات الحساب بنجاح!")

    # ── عرض بيانات الحساب ──
    if "account_user" in st.session_state:
        render_user_card(st.session_state["account_user"])
        if "account_tweet" in st.session_state:
            with st.expander("📌 آخر تغريدة مجلوبة"):
                render_tweet_card(st.session_state["account_tweet"])

    # ═══════════════════════════════════════════════
    #  قسم التحليل الاستخباراتي الشامل للمنشورات
    # ═══════════════════════════════════════════════
    st.divider()
    st.subheader("🔴 التحليل الاستخباراتي الشامل للمنشورات")

    st.markdown("""
    <div style="background:#1a1a2e;border-radius:10px;padding:14px;
                border-right:4px solid #e94560;direction:rtl;
                margin-bottom:16px;font-size:0.9rem;color:#ccc;">
        <strong style="color:#e94560;">كيفية الاستخدام:</strong><br>
        <strong>الطريقة 1 (تلقائية):</strong> أدخل مفتاح TwitterAPI.io
        واضغط "جلب آخر 500 منشور" – سيجلب المنشورات تلقائياً.<br>
        <strong>الطريقة 2 (يدوية):</strong> الصق نصوص المنشورات مباشرة في الحقل أدناه.
    </div>
    """, unsafe_allow_html=True)

    # ── إعدادات TwitterAPI.io ──
    with st.expander("⚙️ إعدادات الجلب التلقائي (TwitterAPI.io)"):
        st.markdown("""
        <div style="direction:rtl;font-size:0.85rem;color:#aaa;margin-bottom:8px;">
            احصل على مفتاح مجاني من
            <a href="https://twitterapi.io" target="_blank"
               style="color:#1da1f2;">twitterapi.io</a>
            (رصيد مجاني $0.1 عند التسجيل، التكلفة ~$0.075 لكل 500 تغريدة)
        </div>
        """, unsafe_allow_html=True)

        tapi_key = st.text_input(
            "🔑 مفتاح TwitterAPI.io",
            type="password",
            key="twitterapi_key",
            placeholder="أدخل المفتاح هنا للجلب التلقائي",
        )
        tweet_limit = st.slider(
            "عدد المنشورات للتحليل",
            min_value=50, max_value=500, value=200, step=50,
            key="tweet_limit_slider",
        )

        if st.button("📥 جلب المنشورات تلقائياً", key="btn_fetch_tweets"):
            if not tapi_key:
                st.error("يرجى إدخال مفتاح TwitterAPI.io.")
            elif "account_user" not in st.session_state:
                st.error("يرجى جلب بيانات الحساب أولاً.")
            else:
                uname = st.session_state["account_user"].get("screen_name", "")
                tweets_list = fetch_user_tweets_twitterapi(
                    uname, tapi_key, limit=tweet_limit
                )
                if tweets_list:
                    st.session_state["intel_tweets"] = tweets_list
                    st.success(
                        f"✅ تم جلب {len(tweets_list)} منشور بنجاح! "
                        "اضغط على 'بدء التحليل الاستخباراتي' أدناه."
                    )
                else:
                    st.error("❌ لم يتم جلب أي منشور.")

    # ── الإدخال اليدوي ──
    st.markdown(
        '<div style="direction:rtl;font-weight:600;color:#ccc;'
        'margin-bottom:6px;">📋 أو الصق نصوص المنشورات يدوياً:</div>',
        unsafe_allow_html=True,
    )
    manual_text = st.text_area(
        "نصوص المنشورات",
        height=200,
        placeholder=(
            "الصق هنا نصوص التغريدات (كل تغريدة في سطر منفصل أو مفصولة بفاصل)\n"
            "مثال:\n"
            "التغريدة الأولى...\n"
            "---\n"
            "التغريدة الثانية...\n"
        ),
        key="manual_tweets_input",
        label_visibility="collapsed",
    )

    # ── معلومات التغريدات المجلوبة ──
    if "intel_tweets" in st.session_state:
        cnt = len(st.session_state["intel_tweets"])
        st.markdown(
            f'<div style="background:#0a1a2a;border-radius:8px;padding:10px 14px;'
            f'border:1px solid #1da1f2;direction:rtl;color:#1da1f2;'
            f'font-size:0.9rem;margin-bottom:10px;">'
            f'📊 التغريدات المجلوبة جاهزة: <strong>{cnt} منشور</strong></div>',
            unsafe_allow_html=True,
        )

    # ── زر التحليل الاستخباراتي ──
    if st.button(
        "🧠 بدء التحليل الاستخباراتي الشامل",
        key="btn_intel_analyze",
        type="primary",
    ):
        if not api_key:
            st.error("يرجى إدخال مفتاح Gemini API في الشريط الجانبي.")
            return

        # جمع النصوص
        tweets_list: list = []

        if "intel_tweets" in st.session_state:
            tweets_list = st.session_state["intel_tweets"]

        if manual_text.strip():
            # تقسيم النص اليدوي
            manual_lines = [
                ln.strip()
                for ln in re.split(r"\n---\n|\n\n", manual_text)
                if ln.strip()
            ]
            tweets_list = manual_lines + tweets_list

        if not tweets_list:
            # تحليل الوصف فقط إن لم تتوفر منشورات
            profile = st.session_state.get("account_user", {})
            desc    = profile.get("description", "")
            if desc:
                tweets_list = [desc]
                st.info(
                    "ℹ️ لا تتوفر منشورات. سيتم التحليل بناءً على وصف الحساب فقط."
                )
            else:
                st.error(
                    "❌ لا توجد منشورات للتحليل.\n\n"
                    "استخدم الجلب التلقائي أو الصق المنشورات يدوياً."
                )
                return

        # 1) الفحص الكلمي
        all_text = "\n".join(tweets_list)
        found_kw = scan_keywords(all_text)

        st.markdown(
            '<div style="direction:rtl;font-weight:700;color:#e94560;'
            'font-size:1rem;margin:12px 0 8px;">📊 نتائج الفحص الكلمي:</div>',
            unsafe_allow_html=True,
        )
        render_keyword_scan(found_kw, len(tweets_list))

        # 2) التحليل بالذكاء الاصطناعي
        with st.spinner("🤖 جاري التحليل الاستخباراتي بالذكاء الاصطناعي..."):
            profile_data = st.session_state.get("account_user", {})
            report = generate_intel_summary(
                tweets_list, profile_data,
                found_kw, api_key, model_name
            )

        if report:
            st.session_state["intel_report"] = report

            # عرض التقرير
            st.markdown(f"""
            <div style="background:linear-gradient(135deg,#1a0510,#2a0a18);
                        border-radius:16px;padding:24px;
                        border:2px solid #b71c1c;
                        box-shadow:0 4px 24px rgba(183,28,28,0.4);
                        direction:rtl;margin-top:16px;line-height:1.9;">
                <h3 style="color:#f44336;margin-bottom:16px;font-size:1.2rem;">
                    🔴 التقرير الاستخباراتي الشامل
                </h3>
                {nl_to_br(report)}
            </div>
            """, unsafe_allow_html=True)

    # ── عرض التقرير المحفوظ ──
    elif "intel_report" in st.session_state:
        st.markdown(f"""
        <div style="background:linear-gradient(135deg,#1a0510,#2a0a18);
                    border-radius:16px;padding:24px;
                    border:2px solid #b71c1c;
                    box-shadow:0 4px 24px rgba(183,28,28,0.4);
                    direction:rtl;margin-top:16px;line-height:1.9;">
            <h3 style="color:#f44336;margin-bottom:16px;font-size:1.2rem;">
                🔴 التقرير الاستخباراتي الشامل
            </h3>
            {nl_to_br(st.session_state['intel_report'])}
        </div>
        """, unsafe_allow_html=True)

    # ════════════════════════════════════
    #  قسم تحليل الصورة
    # ════════════════════════════════════
    st.divider()
    st.subheader("🖼️ تحليل صورة الملف الشخصي")

    images_b64: list = []

    if uploaded_img:
        img = Image.open(uploaded_img)
        b64 = image_to_base64(img)
        images_b64.append(b64)
        st.image(img, caption="الصورة المرفوعة", width=280)
    elif "account_user" in st.session_state:
        avatar = st.session_state["account_user"].get("avatar_url", "")
        if avatar:
            avatar_hq = avatar.replace("_normal", "_400x400")
            b64 = download_image_b64(avatar_hq) or download_image_b64(avatar)
            if b64:
                images_b64.append(b64)
                st.image(base64_to_bytesio(b64),
                         caption="صورة الملف الشخصي", width=200)

    if images_b64:
        point = st.selectbox(
            "نقطة تحليل الصورة",
            list(IMAGE_ANALYSIS_POINTS.keys()),
            key="account_img_point",
        )
        if st.button("🔍 تحليل الصورة", key="btn_account_img_analyze"):
            if not api_key:
                st.error("يرجى إدخال مفتاح Gemini API.")
                return
            prompt = IMAGE_ANALYSIS_POINTS[point]
            with st.spinner("🤖 جاري التحليل..."):
                result = gemini_with_images(prompt, images_b64, api_key, model_name)
            if result:
                st.session_state["img_analysis"] = result
                st.session_state["images_b64"]   = images_b64
                st.markdown(
                    '<div style="background:#1a1a2e;border-radius:12px;'
                    'padding:20px;border:1px solid #e94560;direction:rtl;'
                    'line-height:1.8;">'
                    + nl_to_br(result) + "</div>",
                    unsafe_allow_html=True,
                )

    render_export_buttons(
        st.session_state.get("account_user", {}).get("screen_name", "account")
    )


# ════════════════════════════════════════════════════
#  تبويب تحليل التغريدة
# ════════════════════════════════════════════════════

def tweet_tab(api_key: str, model_name: str):
    st.header("📌 تحليل تغريدة")

    col1, col2 = st.columns([3, 1])
    with col1:
        raw_url = st.text_input(
            "رابط التغريدة أو معرّفها الرقمي",
            placeholder="https://x.com/user/status/1234567890",
            key="tweet_url_input",
        )
    with col2:
        uploaded_imgs = st.file_uploader(
            "صور إضافية",
            type=["jpg", "png", "webp"],
            accept_multiple_files=True,
            key="tweet_imgs_upload",
        )

    if st.button("🔎 جلب التغريدة", key="btn_fetch_tweet"):
        if not raw_url.strip():
            st.error("يرجى إدخال الرابط أو المعرّف.")
            return

        tweet_id = extract_tweet_id(raw_url)
        if not tweet_id:
            st.error("❌ تعذّر استخراج معرّف التغريدة. تأكد أن الرابط يحتوي على /status/")
            return

        with st.spinner("⏳ جاري جلب التغريدة..."):
            tweet = fetch_fxtwitter_tweet(tweet_id)

        if not tweet:
            st.error("❌ تعذّر جلب التغريدة. تحقق من الرابط.")
            return

        author = tweet.get("author", {})
        td = {
            "المعرّف":   tweet.get("id", ""),
            "النص":      tweet.get("text", ""),
            "المؤلف":    author.get("name", ""),
            "الحساب":    f"@{author.get('screen_name', '')}",
            "التاريخ":   format_date(tweet.get("created_timestamp", 0)),
            "الإعجابات": format_number(tweet.get("likes",    0)),
            "إعادة نشر": format_number(tweet.get("retweets", 0)),
            "الردود":    format_number(tweet.get("replies",  0)),
            "المشاهدات": format_number(tweet.get("views",    0)),
            "اللغة":     tweet.get("lang", ""),
        }
        st.session_state["tweet_data"]    = td
        st.session_state["tweet_obj"]     = tweet
        st.session_state["text_analysis"] = ""
        st.session_state["exec_summary"]  = ""
        st.session_state.pop("img_analysis", None)
        st.success("✅ تم جلب التغريدة بنجاح!")

    if "tweet_obj" not in st.session_state:
        return

    render_tweet_card(st.session_state["tweet_obj"])

    # ── صور التغريدة ──
    tweet_obj      = st.session_state["tweet_obj"]
    tweet_imgs_b64: list = []

    media = tweet_obj.get("media", {})
    if isinstance(media, dict):
        for photo in media.get("photos", []):
            url = photo.get("url") or photo.get("thumbnail_url")
            if url:
                b64 = download_image_b64(url)
                if b64:
                    tweet_imgs_b64.append(b64)
                    st.image(base64_to_bytesio(b64),
                             caption="صورة من التغريدة", width=500)

    for uf in (uploaded_imgs or []):
        img = Image.open(uf)
        b64 = image_to_base64(img)
        tweet_imgs_b64.append(b64)
        st.image(img, caption=f"صورة مرفوعة: {uf.name}", width=500)

    if tweet_imgs_b64:
        st.session_state["images_b64"] = tweet_imgs_b64

    # ── تحليل الصورة ──
    if tweet_imgs_b64:
        st.divider()
        st.subheader("🖼️ تحليل الصورة")
        point = st.selectbox(
            "نقطة التحليل",
            list(IMAGE_ANALYSIS_POINTS.keys()),
            key="tweet_img_point",
        )
        if st.button("🔍 تحليل الصورة", key="btn_tweet_img_analyze"):
            if not api_key:
                st.error("يرجى إدخال مفتاح Gemini API.")
                return
            with st.spinner("🤖 جاري تحليل الصورة..."):
                result = gemini_with_images(
                    IMAGE_ANALYSIS_POINTS[point],
                    tweet_imgs_b64, api_key, model_name
                )
            if result:
                st.session_state["img_analysis"] = result
                st.markdown(
                    '<div style="background:#1a1a2e;border-radius:12px;'
                    'padding:20px;border:1px solid #1da1f2;direction:rtl;'
                    'line-height:1.8;">'
                    + nl_to_br(result) + "</div>",
                    unsafe_allow_html=True,
                )
    elif st.session_state.get("img_analysis"):
        st.markdown(
            '<div style="background:#1a1a2e;border-radius:12px;'
            'padding:20px;border:1px solid #1da1f2;direction:rtl;line-height:1.8;">'
            + nl_to_br(st.session_state["img_analysis"]) + "</div>",
            unsafe_allow_html=True,
        )

    # ── تحليل النص ──
    st.divider()
    st.subheader("📝 تحليل نص التغريدة")
    tweet_text = st.session_state["tweet_data"].get("النص", "")

    if tweet_text:
        st.markdown(
            '<div style="background:#0d1b2a;border-radius:8px;padding:14px;'
            'border-right:4px solid #1da1f2;direction:rtl;color:#ccc;'
            'font-size:0.95rem;line-height:1.7;margin-bottom:12px;">'
            + nl_to_br(tweet_text) + "</div>",
            unsafe_allow_html=True,
        )

    if st.button("📝 تحليل النص بالذكاء الاصطناعي", key="btn_text_analysis"):
        if not api_key:
            st.error("يرجى إدخال مفتاح Gemini API.")
            return
        if not tweet_text:
            st.error("النص فارغ.")
            return
        prompt = (
            "قم بتحليل استخباراتي شامل لهذه التغريدة:\n\n"
            f'"{tweet_text}"\n\n'
            "اشمل:\n1. المحتوى الرئيسي\n2. النية والهدف\n"
            "3. نوع الخطاب\n4. الكلمات المفتاحية\n"
            "5. مستوى الخطورة (1-10)\n6. تقييم المصداقية\n"
            "7. السياق والخلفية"
        )
        with st.spinner("🤖 جاري التحليل..."):
            result = gemini_text(prompt, api_key, model_name)
        if result:
            st.session_state["text_analysis"] = result
            st.markdown(
                '<div style="background:#1a1a2e;border-radius:12px;'
                'padding:20px;border:1px solid #1da1f2;direction:rtl;'
                'line-height:1.8;">'
                + nl_to_br(result) + "</div>",
                unsafe_allow_html=True,
            )
    elif st.session_state.get("text_analysis"):
        st.markdown(
            '<div style="background:#1a1a2e;border-radius:12px;'
            'padding:20px;border:1px solid #1da1f2;direction:rtl;line-height:1.8;">'
            + nl_to_br(st.session_state["text_analysis"]) + "</div>",
            unsafe_allow_html=True,
        )

    # ── الملخص التنفيذي ──
    st.divider()
    st.subheader("🧠 الملخص التنفيذي")

    if st.button("🧠 توليد الملخص التنفيذي الشامل", key="btn_exec_summary"):
        if not api_key:
            st.error("يرجى إدخال مفتاح Gemini API.")
            return
        img_a  = st.session_state.get("img_analysis",  "")
        text_a = st.session_state.get("text_analysis", "")
        if not (img_a or text_a):
            st.warning("يرجى إجراء تحليل الصورة أو النص أولاً.")
            return
        prompt_es = (
            "بناءً على التحليلات التالية، أنشئ ملخصاً تنفيذياً استخباراتياً "
            "احترافياً يشمل:\n"
            "1. 📌 النتائج الرئيسية\n2. 🔍 المؤشرات الاستخباراتية\n"
            "3. ⚠️ مستوى الخطورة (1-10)\n4. ✅ تقييم المصداقية\n"
            "5. 🌐 السياق\n6. ⚡ التوصيات الفورية\n"
            "7. 📋 الاستراتيجية المقترحة\n\n"
            f"--- تحليل الصور ---\n{img_a}\n\n"
            f"--- تحليل النص ---\n{text_a}"
        )
        with st.spinner("🧠 جاري توليد الملخص..."):
            summary = gemini_text(prompt_es, api_key, model_name)
        if summary:
            st.session_state["exec_summary"] = summary
            st.markdown(f"""
            <div style="background:linear-gradient(135deg,#1a0a10,#2a1020);
                        border-radius:16px;padding:24px;
                        border:2px solid #e94560;
                        box-shadow:0 4px 24px rgba(233,69,96,0.35);
                        direction:rtl;margin-top:12px;line-height:1.9;">
                <h3 style="color:#e94560;margin-bottom:16px;font-size:1.2rem;">
                    📋 الملخص التنفيذي
                </h3>
                {nl_to_br(summary)}
            </div>
            """, unsafe_allow_html=True)

    elif st.session_state.get("exec_summary"):
        st.markdown(f"""
        <div style="background:linear-gradient(135deg,#1a0a10,#2a1020);
                    border-radius:16px;padding:24px;
                    border:2px solid #e94560;
                    box-shadow:0 4px 24px rgba(233,69,96,0.35);
                    direction:rtl;margin-top:12px;line-height:1.9;">
            <h3 style="color:#e94560;margin-bottom:16px;font-size:1.2rem;">
                📋 الملخص التنفيذي
            </h3>
            {nl_to_br(st.session_state['exec_summary'])}
        </div>
        """, unsafe_allow_html=True)

    render_export_buttons(
        st.session_state.get("tweet_data", {}).get("الحساب", "tweet")
    )


# ════════════════════════════════════════════════════
#  الشريط الجانبي
# ════════════════════════════════════════════════════

def render_sidebar():
    with st.sidebar:
        st.markdown(
            f'<div style="text-align:center;padding:10px 0;">'
            f'<span style="font-size:2rem;">{PAGE_ICON}</span>'
            f'<h3 style="color:#e94560;margin:4px 0;">{PAGE_TITLE}</h3>'
            f'<span style="color:#666;font-size:0.85rem;">الإصدار: {VERSION}</span>'
            f'</div>',
            unsafe_allow_html=True,
        )
        st.divider()

        api_key = st.text_input(
            "🔑 مفتاح Gemini API",
            type="password",
            help="احصل على مفتاحك من: https://aistudio.google.com/apikey",
            key="gemini_api_key",
        )

        model_label = st.selectbox(
            "🤖 النموذج",
            list(GEMINI_MODELS.keys()),
            key="model_selector",
        )
        model_name = GEMINI_MODELS[model_label]

        st.divider()
        st.markdown("#### 📡 مصادر البيانات")
        st.markdown(
            '<div style="background:#0a2a0a;border-radius:8px;padding:8px 12px;'
            'color:#4caf50;font-size:0.85rem;">✅ FxTwitter API — نشط</div>',
            unsafe_allow_html=True,
        )
        st.markdown(
            '<div style="background:#1a1a0a;border-radius:8px;padding:8px 12px;'
            'color:#ff9800;font-size:0.85rem;margin-top:6px;">'
            '⚡ TwitterAPI.io — اختياري (للجلب التلقائي)</div>',
            unsafe_allow_html=True,
        )
        st.markdown(
            '<div style="background:#2a0a0a;border-radius:8px;padding:8px 12px;'
            'color:#f44336;font-size:0.85rem;margin-top:6px;">'
            '❌ Nitter — مغلق (2024)</div>',
            unsafe_allow_html=True,
        )

        st.divider()
        st.markdown("#### 📖 طريقة الاستخدام")
        st.markdown("""
        **تحليل حساب:**
        - أدخل اسم المستخدم أو رابط الحساب
        - للتحليل الشامل: أضف مفتاح TwitterAPI.io

        **تحليل تغريدة:**
        - أدخل رابط التغريدة كاملاً

        **التحليل الاستخباراتي:**
        - جلب تلقائي أو لصق يدوي للمنشورات
        - يحتاج مفتاح Gemini API
        """)

    return api_key, model_name


# ════════════════════════════════════════════════════
#  الدالة الرئيسية
# ════════════════════════════════════════════════════

def main():
    st.set_page_config(
        page_title=PAGE_TITLE,
        page_icon=PAGE_ICON,
        layout="wide",
        initial_sidebar_state="expanded",
    )

    st.markdown("""
    <style>
    body, .stApp {
        background-color: #0d1b2a !important;
        color: #e0e0e0 !important;
    }
    .stApp, .main, [data-testid="stAppViewContainer"],
    [data-testid="stVerticalBlock"],
    .stMarkdown, .stText,
    .stAlert, .stWarning, .stError, .stSuccess, .stInfo,
    h1, h2, h3, h4, h5, h6, p, li, span, div {
        direction: rtl !important;
        text-align: right !important;
    }
    .stTextInput label, .stSelectbox label,
    .stFileUploader label, .stNumberInput label,
    .stCheckbox label, .stRadio label,
    .stTextArea label, .stSlider label {
        direction: rtl !important;
        text-align: right !important;
        font-weight: 600 !important;
        color: #bbb !important;
        width: 100%;
        display: block;
    }
    .stTextInput input, .stNumberInput input {
        direction: rtl !important;
        text-align: right !important;
        background-color: #1a2a3a !important;
        color: #fff !important;
        border: 1px solid #2a3a4a !important;
        border-radius: 8px !important;
        padding: 10px 14px !important;
    }
    .stTextInput input:focus, .stNumberInput input:focus {
        border-color: #e94560 !important;
        box-shadow: 0 0 0 2px rgba(233,69,96,0.25) !important;
        outline: none !important;
    }
    .stTextInput input[type="password"] {
        direction: ltr !important;
        text-align: left !important;
    }
    textarea {
        direction: rtl !important;
        text-align: right !important;
        background-color: #1a2a3a !important;
        color: #fff !important;
        border: 1px solid #2a3a4a !important;
        border-radius: 8px !important;
    }
    textarea:focus {
        border-color: #e94560 !important;
        box-shadow: 0 0 0 2px rgba(233,69,96,0.25) !important;
    }
    .stSelectbox > div > div {
        background-color: #1a2a3a !important;
        direction: rtl !important;
        text-align: right !important;
        border-radius: 8px !important;
    }
    .stButton > button {
        background: linear-gradient(135deg, #e94560, #c0392b) !important;
        color: white !important;
        border: none !important;
        border-radius: 8px !important;
        font-weight: 600 !important;
        padding: 10px 22px !important;
        width: 100%;
        transition: all 0.2s ease;
        font-size: 0.95rem !important;
    }
    .stButton > button:hover {
        opacity: 0.88 !important;
        transform: translateY(-1px) !important;
        box-shadow: 0 4px 12px rgba(233,69,96,0.4) !important;
    }
    .stDownloadButton > button {
        background: linear-gradient(135deg, #1da1f2, #0d7abf) !important;
        color: white !important;
        border: none !important;
        border-radius: 8px !important;
        font-weight: 600 !important;
        width: 100%;
    }
    .stTabs [data-baseweb="tab"] {
        font-size: 1rem !important;
        font-weight: 600 !important;
        padding: 12px 28px !important;
    }
    .stTabs [aria-selected="true"] {
        border-bottom: 3px solid #e94560 !important;
        color: #e94560 !important;
    }
    [data-testid="stAlert"] > div,
    .stAlert > div, .stWarning > div,
    .stError > div, .stSuccess > div {
        direction: rtl !important;
        text-align: right !important;
    }
    .streamlit-expanderHeader {
        direction: rtl !important;
        text-align: right !important;
        font-weight: 600 !important;
    }
    [data-testid="stSidebar"],
    [data-testid="stSidebar"] * {
        direction: rtl !important;
        text-align: right !important;
    }
    [data-testid="stSidebar"] {
        background-color: #111827 !important;
    }
    .stSlider > div { direction: rtl !important; }
    hr { border-color: #2a3a4a !important; }
    [data-testid="stFileUploader"] { direction: rtl !important; }
    [data-testid="stFileUploader"] > div {
        direction: rtl !important;
        text-align: right !important;
    }
    .stSpinner > div {
        direction: rtl !important;
        text-align: right !important;
    }
    ::-webkit-scrollbar { width: 5px; height: 5px; }
    ::-webkit-scrollbar-track { background: #0d1b2a; }
    ::-webkit-scrollbar-thumb { background: #e94560; border-radius: 3px; }
    </style>
    """, unsafe_allow_html=True)

    api_key, model_name = render_sidebar()

    tab1, tab2 = st.tabs(["🔍  تحليل حساب X", "📌  تحليل تغريدة"])
    with tab1:
        account_tab(api_key, model_name)
    with tab2:
        tweet_tab(api_key, model_name)


if __name__ == "__main__":
    main()
